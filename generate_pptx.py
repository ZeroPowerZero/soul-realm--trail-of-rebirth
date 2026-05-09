import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def main():
    prs = Presentation()
    
    # Use standard 16:9 aspect ratio if possible, otherwise stick to default (4:3)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]

    # Helper function to style titles
    def style_title(shape, text, font_size=40):
        shape.text = text
        for p in shape.text_frame.paragraphs:
            p.font.name = 'Arial'
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(108, 92, 231) # Purple-ish
            
    def add_bullet_slide(prs, title, bullets):
        slide = prs.slides.add_slide(content_layout)
        style_title(slide.shapes.title, title)
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            p.space_after = Pt(14)
        return slide

    def add_flow_diagram(slide, steps, top_in, left_in=1.0, width_in=2.0, gap_in=0.5):
        current_left = left_in
        for i, step in enumerate(steps):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(current_left), Inches(top_in), Inches(width_in), Inches(1.0)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(76, 29, 149)
            shape.line.color.rgb = RGBColor(139, 92, 246)
            tf = shape.text_frame
            tf.text = step
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            
            current_left += width_in + gap_in
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, 
                    Inches(current_left - gap_in + 0.1), Inches(top_in + 0.35), Inches(gap_in - 0.2), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(236, 72, 153) # Pink
                arrow.line.color.rgb = RGBColor(236, 72, 153)

    # 1. Title
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Soul Realm: Trail of Rebirth"
    subtitle.text = "System Implementation, Architecture, and Final Vision\n\nGodot 4.6 | GDScript | Jolt Physics"
    for p in title.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(139, 92, 246)

    # 2. Final Vision
    slide = add_bullet_slide(prs, "1. The Final Vision (End Result)", [
        "A seamless, high-fidelity roguelite experience.",
        "Procedural Tartarus Dungeons: Hades-style room-to-room progression.",
        "Complex Multi-Phase Bosses: Dynamic behavior based on FSM states.",
        "Limitless Magic Scaling: Spells evolve from projectiles to AoE anomalies.",
        "Metaprogression: Permanent upgrades between runs."
    ])
    add_flow_diagram(slide, ["Hub World", "Dungeon Gen", "Combat Loop", "Boss Fight"], 5.0, left_in=1.5, width_in=2.0)

    # 3. Component Architecture
    slide = add_bullet_slide(prs, "2. Architecture: Components over Inheritance", [
        "Problem: Deep inheritance trees ('god scripts') are rigid and prone to bugs.",
        "Solution: Composition. Behaviors are modular nodes attached to entities.",
        "HealthComponent: Manages HP, emits death signals.",
        "ManaComponent: Handles spell casting economy and regeneration.",
        "UpgradeComponent: Tracks and applies stat modifiers."
    ])

    # 4. Visual Example: Components
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1)).text_frame.text = "Visual Example: Building a Boss"
    style_title(slide.shapes[0], "3. End Result: The Modular Entity")
    
    # Draw central entity
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3), Inches(3), Inches(3))
    base.text = "Boss Node\n(KinematicBody3D)"
    
    comps = ["HealthComponent\n(HP: 5000)", "SpellController\n(Data-Driven Magic)", "StateMachine\n(Phase Logic)"]
    for i, c in enumerate(comps):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1 + i*4), Inches(1.5), Inches(2.5), Inches(1))
        shape.text = c
        # Draw connector
        from pptx.enum.shapes import MSO_CONNECTOR
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.25 + i*4), Inches(2.5), Inches(6.5), Inches(3))

    # 5. Spell Recognition Pipeline
    slide = add_bullet_slide(prs, "4. Spell Recognition Algorithm", [
        "Based on the $1 Unistroke Recognizer algorithm.",
        "Allows players to draw symbols in real-time to cast spells.",
        "Algorithm focuses on shape, not scale or speed.",
        "Enables skill-based, immersive combat mechanics."
    ])
    add_flow_diagram(slide, ["Capture", "Resample", "Center", "Scale", "Match"], 5.0, left_in=0.5, width_in=2.0, gap_in=0.5)

    # 6. Gesture Logic
    slide = add_bullet_slide(prs, "5. Recognition Implementation (Logic)", [
        "Stroke Capture: Mouse/touch inputs are captured via Godot _gui_input.",
        "Interpolation: If the mouse moves too fast, gaps are mathematically filled.",
        "Resampling: Path is divided into exactly 64 equidistant points.",
        "Translation: The shape's centroid is moved to (0,0).",
        "Scaling: The shape is bounded within a 250x250 square."
    ])

    # 7. Data Driven Spells
    slide = add_bullet_slide(prs, "6. Data-Driven Spell Mechanics", [
        "Separation of concerns: Data vs. Behavior.",
        "SpellData (Resource): Static config (Name, Element, Cost, Scene).",
        "SpellDriver (Resource): Runtime state (Level, Damage multiplier).",
        "SpellController (Node): The 3D spawner handling orientation and instantiation.",
        "End Result: Enemies and players use the exact same logic to cast."
    ])

    # 8. Fireball Example
    slide = add_bullet_slide(prs, "7. Visual Example: Fireball Scaling", [
        "Level 1: Single projectile, basic damage.",
        "Level 2: Penetration - passes through 2 enemies.",
        "Level 3: AoE Explosion - erupts on impact or lifetime expiry.",
        "Level 4: Multi-Cast - automatically spawns 3 fireballs sequentially.",
        "Logic: Code checks level flags inside _physics_process to trigger extra effects."
    ])

    # 9. FSM
    slide = add_bullet_slide(prs, "8. Finite State Machine (FSM)", [
        "The Brain of the Player and Bosses.",
        "Prevents complex input conflicts (e.g., casting while dead).",
        "States: Idle, Move, Spell, Dash, Fall, Hit, Death.",
        "Each state is an isolated script handling its own _process and _input."
    ])
    add_flow_diagram(slide, ["Idle", "Move", "Spell", "Dash"], 5.5, left_in=1.5, width_in=2.0)

    # 10. Dash Logic
    slide = add_bullet_slide(prs, "9. FSM Logic Example: The Dash", [
        "Dash State requires physics validation to prevent clipping.",
        "1. Calculate intended dash vector.",
        "2. Cast a PhysicsRayQueryParameters3D raycast into the future.",
        "3. If a wall is detected, stop 1.0 unit before the wall.",
        "4. Execute smooth movement tween to the validated destination."
    ])

    # 11. Enemy AI
    slide = add_bullet_slide(prs, "10. Combat AI: The MageBase Framework", [
        "Abstract class handling shared enemy behavior.",
        "Auto-generates HealthComponent and NavigationAgent3D.",
        "Handles dynamic floating 3D health bars using Viewports.",
        "Implements Line of Sight (LoS) raycasting to prevent shooting through walls.",
        "To create a new enemy, simply inherit MageBase and assign spells."
    ])

    # 12. CounterMage AI
    slide = add_bullet_slide(prs, "11. Reactive AI: The CounterMage", [
        "End Result: A dynamic chess match where enemies react to the player.",
        "Logic: CounterMage secretly listens to the Player's 'drawing_state_changed' signal.",
        "Action: If the player starts drawing and is in Line of Sight...",
        "Response: The Mage halts movement and casts a 0.4s fast interrupt spell.",
        "Forces the player to use environmental cover strategically."
    ])

    # 13. Boss Phases
    slide = add_bullet_slide(prs, "12. End Result: Multi-Phase Bosses", [
        "Combining FSM + Components creates epic encounters.",
        "Phase 1 (100%-50% HP): Standard ranged AI patterns.",
        "Trigger: HealthComponent 'health_changed' signal detects <50% HP.",
        "Transition: StateMachine forces 'PhaseTwoState'.",
        "Phase 2: UpgradeComponent applies +200% Cast Speed buff.",
        "New SpellData arrays are loaded instantly."
    ])

    # 14. Dungeons
    slide = add_bullet_slide(prs, "13. Procedural Dungeon Generation", [
        "Hades-style room-to-room progression.",
        "RoomController: Pre-built scenes with enemy spawns.",
        "DoorController: Handles transitions. Doors lock upon entry.",
        "Clear Condition: RoomController detects all enemy deaths.",
        "Doors unlock, map manager assigns next room type (Combat, Elite, Boss)."
    ])

    # 15. Roguelite Loop
    slide = add_bullet_slide(prs, "14. The Roguelite Progression Loop", [
        "Core gameplay cycle rewarding combat success.",
        "1. Enemy Death triggers GameManager.add_xp().",
        "2. Dynamic Thresholds: Required XP scales exponentially (x1.2 per level).",
        "3. Game pauses, Upgrade UI presented.",
        "4. UpgradeManager applies chosen buff via UpgradeComponent."
    ])
    add_flow_diagram(slide, ["Combat", "Gain XP", "Level Up", "Select Buff"], 5.5, left_in=1.5, width_in=2.0)

    # 16. RNG Logic
    slide = add_bullet_slide(prs, "15. Dynamic Upgrade Selection Logic", [
        "Weighted Random Number Generation (RNG).",
        "Rarity Tiers: Common (High weight), Legendary (Low weight).",
        "Logic filter removes upgrades the player has maxed out.",
        "Selects 3 unique options by rolling against the total pool weight.",
        "Stat upgrades stack recursively (e.g., +20 Max Health per level)."
    ])

    # 17. Aim Assist
    slide = add_bullet_slide(prs, "16. Game Feel: Aim Assist Math", [
        "Makes spellcasting fluid in a 3D environment.",
        "Target Selection: Uses Dot Product of forward vector and enemy direction.",
        "Combined Score = Angular Accuracy - (Distance / MaxRange * 0.5).",
        "Homing: Projectiles use Lerp to gently steer towards the target over time.",
        "Prevents instant snapping, maintaining visual realism."
    ])

    # 18. Camera Juice
    slide = add_bullet_slide(prs, "17. Game Feel: Camera Juice", [
        "Walk Bob: Sinusoidal math (sin(time)) applied to camera Y offset.",
        "FOV Punch: Tweening FOV from 75 to 85 during Dash.",
        "Camera Tilt: Rotating Z-axis inversely to horizontal mouse velocity.",
        "Camera Shake: Multi-step random offset tween on damage taken.",
        "Landing Bob: Downward dip upon ground contact."
    ])

    # 19. Summary
    slide = add_bullet_slide(prs, "18. Summary: Systems Ready for Content", [
        "The technical foundation is complete and highly scalable.",
        "Component Architecture = Zero technical debt when adding entities.",
        "FSM = Zero input conflicts during intense combat.",
        "Data-Driven Spells = Infinite magic variety with zero new scripts.",
        "Next Steps: Content generation (boss assets, room layouts)."
    ])

    # 20. Q&A
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = "Questions & Answers\n\nGodot 4.6 Development"
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(139, 92, 246)

    prs.save(r'c:\Users\sharm\OneDrive\Desktop\soul_realm_presentation.pptx')

if __name__ == '__main__':
    main()
